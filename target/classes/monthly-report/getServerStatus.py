# -*- coding: utf-8 -*-
import os
import sys
from datetime import datetime

# 패키지 import 시도
try:
    from elasticsearch import Elasticsearch
    ELASTICSEARCH_AVAILABLE = True
except ImportError:
    ELASTICSEARCH_AVAILABLE = False
    print("elasticsearch 패키지가 없습니다. Elasticsearch 기능을 건너뜁니다.")

try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False
    print("requests 패키지가 없습니다. HTTP 요청 기능을 건너뜁니다.")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("python-pptx 패키지가 없습니다. PowerPoint 기능을 건너뜁니다.")

import re
from collections import defaultdict
import math
import urllib.parse
import json

# 결과를 저장할 리스트
results = []

def save_result(section, data):
    """결과를 리스트에 저장하는 함수"""
    results.append(f"{section}")
    if isinstance(data, list):
        for item in data:
            results.append(item)
    else:
        results.append(data)
    results.append("")  # 빈 줄 추가

def analyze_powerpoint(pptx_path):
    """PowerPoint 파일의 구조를 분석하는 함수"""
    try:
        prs = Presentation(pptx_path)
        print(f"\nPowerPoint 파일 분석: {pptx_path}")
        print(f"총 슬라이드 수: {len(prs.slides)}")
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n--- 슬라이드 {slide_idx + 1} ---")
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text.strip()
                    if text:
                        print(f"텍스트박스 {shape_idx + 1}: {text[:100]}...")
                elif hasattr(shape, "table"):
                    table = shape.table
                    print(f"표 {shape_idx + 1}: {table.rows}행 x {table.columns}열")
                    # 표 내용 출력
                    for row_idx, row in enumerate(table.rows):
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        if any(row_text):  # 빈 행이 아닌 경우만
                            print(f"  행 {row_idx + 1}: {' | '.join(row_text)}")
    except Exception as e:
        print(f"PowerPoint 파일 분석 중 오류: {e}")

def update_template_powerpoint(pptx_path, hdd_data, cpu_data, memory_data, docs_data):
    """템플릿 PowerPoint 파일의 플레이스홀더를 실제 데이터로 교체하는 함수"""
    try:
        from pptx import Presentation
        import re
        
        # PowerPoint 파일 열기
        prs = Presentation(pptx_path)
        
        # 데이터를 딕셔너리로 변환
        data_dict = {}
        
        # 날짜 데이터 추가 (mm월 dd일 형식)
        current_date = datetime.now()
        date_str = f"{current_date.month:02d}월 {current_date.day:02d}일"
        data_dict["date"] = date_str
        
        # HDD 데이터 파싱 (RC10 876GB/1741GB (50%) → 876GB/1741GB (50%))
        for line in hdd_data:
            if "RC" in line:
                rack = line.split()[0]
                # RC10 876GB/1741GB (50%) → 876GB/1741GB (50%)
                hdd_value = line.replace(rack + " ", "")
                data_dict[f"{rack.lower()}_hdd"] = hdd_value
        
        # CPU 데이터 파싱 ({group="RC10"} 2% → 2%)
        for line in cpu_data:
            if "RC" in line:
                rack = line.split('"')[1] if '"' in line else line.split()[0]
                # {group="RC10"} 2% → 2%
                cpu_value = line.split()[-1]  # 마지막 부분만 (2%)
                data_dict[f"{rack.lower()}_cpu"] = cpu_value
        
        # Memory 데이터 파싱 ({group="RC10"} 66% → 66%)
        for line in memory_data:
            if "RC" in line:
                rack = line.split('"')[1] if '"' in line else line.split()[0]
                # {group="RC10"} 66% → 66%
                memory_value = line.split()[-1]  # 마지막 부분만 (66%)
                data_dict[f"{rack.lower()}_memory"] = memory_value
        
        # 문서 수 데이터 파싱 (RC10 총 문서 수 (Primary Shard 기준): 2,585,773,384개 → 2,585,773,384)
        for line in docs_data:
            if "RC" in line:
                rack = line.split()[0]
                # RC10 총 문서 수 (Primary Shard 기준): 2,585,773,384개 → 2,585,773,384
                docs_value = line.split(":")[-1].strip().replace("개", "")
                data_dict[f"{rack.lower()}_docs"] = docs_value
        
        print(f"파싱된 데이터: {data_dict}")
        
        # 각 슬라이드에서 플레이스홀더 찾아서 교체
        for slide_idx, slide in enumerate(prs.slides):
            print(f"\n슬라이드 {slide_idx + 1} 처리 중...")
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame"):
                    # 텍스트 프레임에서 플레이스홀더 찾기
                    original_text = shape.text_frame.text
                    updated_text = original_text
                    
                    # 모든 플레이스홀더 패턴 찾기: {rc10_hdd}, {rc12_cpu} 등
                    placeholders = re.findall(r'\{([^}]+)\}', original_text)
                    
                    if placeholders:
                        print(f"  텍스트박스 {shape_idx + 1}에서 플레이스홀더 발견: {placeholders}")
                        
                        for placeholder in placeholders:
                            if placeholder in data_dict:
                                updated_text = updated_text.replace(f"{{{placeholder}}}", data_dict[placeholder])
                                print(f"    {placeholder} → {data_dict[placeholder]}")
                            else:
                                print(f"    {placeholder}에 대한 데이터를 찾을 수 없습니다.")
                        
                        # 텍스트 업데이트
                        if updated_text != original_text:
                            shape.text_frame.text = updated_text
                            
                            # 폰트 및 정렬 설정
                            for paragraph in shape.text_frame.paragraphs:
                                # 날짜 플레이스홀더는 우측 정렬, 나머지는 가운데 정렬
                                if '{date}' in original_text:
                                    paragraph.alignment = PP_ALIGN.RIGHT
                                else:
                                    paragraph.alignment = PP_ALIGN.CENTER
                                    
                                for run in paragraph.runs:
                                    run.font.name = 'Pretendard Light'
                                    # 날짜 플레이스홀더는 7pt, 나머지는 9pt
                                    if '{date}' in original_text:
                                        run.font.size = Pt(7)
                                    else:
                                        run.font.size = Pt(9)
                            
                            font_size = "7pt" if '{date}' in original_text else "9pt"
                            alignment = "우측" if '{date}' in original_text else "가운데"
                            print(f"    텍스트 업데이트 완료 (폰트: Pretendard Light, 크기: {font_size}, 정렬: {alignment})")
                
                elif hasattr(shape, "table"):
                    # 표에서도 플레이스홀더 찾기
                    table = shape.table
                    print(f"  표 {shape_idx + 1}: {table.rows}행 x {table.columns}열")
                    
                    for row_idx, row in enumerate(table.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            cell_text = cell.text
                            
                            if cell_text and '{' in cell_text:
                                print(f"    셀 [{row_idx+1},{cell_idx+1}]: {cell_text}")
                                
                                # 셀에서 플레이스홀더 찾기
                                placeholders = re.findall(r'\{([^}]+)\}', cell_text)
                                
                                for placeholder in placeholders:
                                    if placeholder in data_dict:
                                        cell.text = cell.text.replace(f"{{{placeholder}}}", data_dict[placeholder])
                                        
                                        # 표 셀의 폰트 및 정렬 설정
                                        for paragraph in cell.text_frame.paragraphs:
                                            # 날짜 플레이스홀더는 우측 정렬, 나머지는 가운데 정렬
                                            if placeholder == 'date':
                                                paragraph.alignment = PP_ALIGN.RIGHT
                                            else:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                
                                            for run in paragraph.runs:
                                                run.font.name = 'Pretendard Light'
                                                # 날짜 플레이스홀더는 7pt, 나머지는 9pt
                                                if placeholder == 'date':
                                                    run.font.size = Pt(7)
                                                else:
                                                    run.font.size = Pt(9)
                                        
                                        print(f"      {placeholder} → {data_dict[placeholder]} (폰트 적용)")
                                    else:
                                        print(f"      {placeholder}에 대한 데이터를 찾을 수 없습니다.")
        
        # 업데이트된 파일 저장 (RnD_월간보고_yyyy년_mm월_현황 형식)
        current_date = datetime.now()
        year = current_date.year
        month = current_date.month
        updated_pptx_path = os.path.join(script_dir, f"RnD_월간보고_{year}년_{month:02d}월_현황.pptx")
        prs.save(updated_pptx_path)
        
        print(f"PowerPoint 템플릿이 업데이트되어 저장되었습니다: {updated_pptx_path}")
        return updated_pptx_path
        
    except Exception as e:
        print(f"PowerPoint 템플릿 업데이트 중 오류 발생: {e}")
        import traceback
        traceback.print_exc()
        return None

print("[HDD 평균 사용량]--------------------------------")

# Elasticsearch 연결
hosts = ['http://203.251.205.71:9200', 'http://203.251.205.72:9200', 'http://203.251.205.73:9200']
es = Elasticsearch(hosts, http_auth=("elastic", "elastic!@#"))

# _cat/allocation API 호출
res = es.cat.allocation(format="json")

# 단위 변환 함수 (TB, GB, MB 등 처리)
def parse_size(s):
    s = s.strip().lower()
    if s.endswith('tb'):
        return float(s.replace('tb', '')) * 1024
    elif s.endswith('gb'):
        return float(s.replace('gb', ''))
    elif s.endswith('mb'):
        return float(s.replace('mb', '')) / 1024
    else:
        return float(s)

# 렉별 디스크 합계 및 노드 수 계산
group_sums = defaultdict(lambda: {"used": 0, "total": 0, "count": 0})

for item in res:
    node = item.get('node')
    if not node:
        continue

    match = re.match(r'(RC\d{2})', node)
    if not match:
        continue

    group = match.group(1)

    try:
        used = parse_size(item.get('disk.used', '0'))
        total = parse_size(item.get('disk.total', '0'))
    except ValueError as e:
        print(f"변환 실패: {item}, 에러: {e}")
        continue

    group_sums[group]["used"] += used
    group_sums[group]["total"] += total
    group_sums[group]["count"] += 1

# 결과 출력 및 저장
hdd_results = []
for group in sorted(group_sums.keys()):
    data = group_sums[group]
    count = data["count"]

    avg_used = math.ceil(data["used"] / count)
    avg_total = math.ceil(data["total"] / count)
    percent = int((avg_used / avg_total) * 100) if avg_total > 0 else 0

    result_line = f"{group} {avg_used}GB/{avg_total}GB ({percent}%)"
    print(result_line)
    hdd_results.append(result_line)

save_result("[HDD 평균 사용량]", hdd_results)

print("[CPU 평균 사용량]--------------------------------")

# Prometheus API 주소
PROMETHEUS_URL = "http://203.251.205.83:9090/api/v1/query"

# 렉별 hostname 정규식
rc_groups = {
    "RC10": "RC10.*",
    "RC12": "RC12.*",
    "RC17": "RC17.*",
    "RC18": "RC18.*"
}

cpu_results = []
for group, hostname_pattern in rc_groups.items():
    # PromQL 개별 쿼리
    promql = (
        f'label_replace('
        f'  avg('
        f'    100 - avg_over_time('
        f'      rate(node_cpu_seconds_total{{mode="idle", hostname=~"{hostname_pattern}", job="Lucy3.0 ES Cluster"}}[5m])'
        f'      [5d:]'
        f'    ) * 100'
        f'  ),'
        f'"group", "{group}", "", ""'
        f')'
    )

    try:
        encoded_query = urllib.parse.quote(promql)
        url = f"{PROMETHEUS_URL}?query={encoded_query}"

        response = requests.get(url)
        if response.status_code == 200:
            data = response.json()
            query_results = data.get("data", {}).get("result", [])

            if not query_results:
                result_line = f'{group} 결과 없음'
                print(result_line)
                cpu_results.append(result_line)
                continue

            for item in query_results:
                value_raw = item.get("value", [])[1]
                try:
                    value_float = float(value_raw)
                    value_rounded = round(value_float)  # 정수 반올림
                    result_line = f'{{group="{group}"}}  {value_rounded}%'
                    print(result_line)
                    cpu_results.append(result_line)
                except ValueError:
                    result_line = f'{{group="{group}"}}  변환 실패: {value_raw}'
                    print(result_line)
                    cpu_results.append(result_line)
        else:
            result_line = f"{group} → HTTP 오류 코드: {response.status_code}"
            print(result_line)
            cpu_results.append(result_line)

    except Exception as e:
        result_line = f"{group} → 예외 발생: {e}"
        print(result_line)
        cpu_results.append(result_line)

save_result("[CPU 평균 사용량]", cpu_results)

print("[Memory 평균 사용량]--------------------------------")

# PromQL 문자열 정의
promql = (
    'label_replace('
    '  avg(('
    '    avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC10.*", job="Lucy3.0 ES Cluster"}[31d]) - '
    '    avg_over_time(node_memory_MemAvailable_bytes{hostname=~"RC10.*", job="Lucy3.0 ES Cluster"}[31d])'
    '  ) / '
    '  avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC10.*", job="Lucy3.0 ES Cluster"}[31d]) * 100), '
    '"group", "RC10", "", ""'
    ') or '
    'label_replace('
    '  avg(('
    '    avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC12.*", job="Lucy3.0 ES Cluster"}[31d]) - '
    '    avg_over_time(node_memory_MemAvailable_bytes{hostname=~"RC12.*", job="Lucy3.0 ES Cluster"}[31d])'
    '  ) / '
    '  avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC12.*", job="Lucy3.0 ES Cluster"}[31d]) * 100), '
    '"group", "RC12", "", ""'
    ') or '
    'label_replace('
    '  avg(('
    '    avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC17.*", job="Lucy3.0 ES Cluster"}[31d]) - '
    '    avg_over_time(node_memory_MemAvailable_bytes{hostname=~"RC17.*", job="Lucy3.0 ES Cluster"}[31d])'
    '  ) / '
    '  avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC17.*", job="Lucy3.0 ES Cluster"}[31d]) * 100), '
    '"group", "RC17", "", ""'
    ') or '
    'label_replace('
    '  avg(('
    '    avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC18.*", job="Lucy3.0 ES Cluster"}[31d]) - '
    '    avg_over_time(node_memory_MemAvailable_bytes{hostname=~"RC18.*", job="Lucy3.0 ES Cluster"}[31d])'
    '  ) / '
    '  avg_over_time(node_memory_MemTotal_bytes{hostname=~"RC18.*", job="Lucy3.0 ES Cluster"}[31d]) * 100), '
    '"group", "RC18", "", ""'
    ')'
)

# 쿼리 파라미터 인코딩
encoded_query = urllib.parse.quote(promql)
full_url = f"{PROMETHEUS_URL}?query={encoded_query}"

memory_results = []
try:
    response = requests.get(full_url)

    if response.status_code == 200:
        data = response.json()
        query_results = data.get("data", {}).get("result", [])

        for item in query_results:
            metric = item.get("metric", {})
            group = metric.get("group", "")
            value_raw = item.get("value", [])[1]

            try:
                value_float = float(value_raw)
                value_rounded = round(value_float)  # 정수로 반올림
                result_line = f'{{group="{group}"}}  {value_rounded}%'
                print(result_line)
                memory_results.append(result_line)
            except ValueError:
                result_line = f'{{group="{group}"}}  변환 실패: {value_raw}'
                print(result_line)
                memory_results.append(result_line)
    else:
        result_line = f"HTTP 오류 코드: {response.status_code}"
        print(result_line)
        memory_results.append(result_line)

except Exception as e:
    result_line = f"예외 발생: {e}"
    print(result_line)
    memory_results.append(result_line)

save_result("[Memory 평균 사용량]", memory_results)

# print("[평균 문서수]--------------------------------")

# # _cat/shards API 호출
# res = es.cat.shards(format="json")

# # 렉별 docs 합계 및 노드 수 계산
# group_docs = defaultdict(int)
# group_node_count = defaultdict(set)  # set으로 중복 제거

# for shard in res:
#     if shard.get('prirep') != 'p':
#         continue  # primary shard만

#     node = shard.get('node')
#     if not node:
#         continue

#     match = re.match(r'(RC\d{2})', node)
#     if not match:
#         continue

#     group = match.group(1)

#     # docs 수 가져오기 (문자열 숫자일 수 있음)
#     docs_str = shard.get('docs', '0').replace(',', '')
#     try:
#         docs = int(docs_str)
#     except ValueError:
#         continue

#     group_docs[group] += docs
#     group_node_count[group].add(node)  # set을 통해 고유 노드 수 저장

# # 결과 출력
# for group in sorted(group_docs.keys()):
#     total_docs = group_docs[group]
#     node_count = len(group_node_count[group])
#     avg_docs = total_docs // node_count if node_count > 0 else 0

#     print(f"{group} 평균 문서 수 (Primary Shard 기준): {avg_docs:,}개")

# print("--------------------------------")
print("[문서 총합]--------------------------------")

# _cat/shards API 호출
res = es.cat.shards(format="json")

# 렉별 docs 합계 계산 (노드 수와 상관없이)
group_docs = defaultdict(int)

for shard in res:
    if shard.get('prirep') != 'p':
        continue  # primary shard만 집계

    node = shard.get('node')
    if not node:
        continue

    match = re.match(r'(RC\d{2})', node)
    if not match:
        continue

    group = match.group(1)

    # docs 수 가져오기
    docs_str = shard.get('docs', '0').replace(',', '')
    try:
        docs = int(docs_str)
    except ValueError:
        continue

    group_docs[group] += docs

# 결과 출력 (그룹별 총합)
docs_results = []
for group in sorted(group_docs.keys()):
    total_docs = group_docs[group]
    result_line = f"{group} 총 문서 수 (Primary Shard 기준): {total_docs:,}"
    print(result_line)
    docs_results.append(result_line)

save_result("[문서 총합]", docs_results)

# 스크립트 디렉토리 설정
script_dir = os.path.dirname(os.path.abspath(__file__))
print(f"스크립트 디렉토리: {script_dir}")
print(f"현재 작업 디렉토리: {os.getcwd()}")

print(f"\n서버 데이터 수집 완료! 바로 PowerPoint 파일을 생성합니다.")

# 템플릿 PowerPoint 파일 업데이트
pptx_template_path = os.path.join(script_dir, "template", "RnD_월간보고_기본템플릿.pptx")
if os.path.exists(pptx_template_path):
    print(f"\n템플릿 PowerPoint 파일 업데이트 중: {pptx_template_path}")
    updated_pptx = update_template_powerpoint(pptx_template_path, hdd_results, cpu_results, memory_results, docs_results)
    if updated_pptx:
        print(f"PowerPoint 템플릿이 성공적으로 업데이트되었습니다!")
        print(f"업데이트된 파일: {updated_pptx}")
    else:
        print("PowerPoint 템플릿 업데이트에 실패했습니다.")
else:
    print(f"PowerPoint 템플릿 파일을 찾을 수 없습니다: {pptx_template_path}")
    print("기본 PowerPoint 파일을 생성합니다...")
    
    # 기본 PowerPoint 파일 생성
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        # 새 프레젠테이션 생성
        prs = Presentation()
        
        # 제목 슬라이드 추가
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Quettai Cluster 시스템 운영 현황"
        subtitle.text = f"{datetime.now().strftime('%Y년 %m월')} 월간 보고서"
        
        # 데이터 슬라이드 추가
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "시스템 현황 요약"
        
        # 텍스트 박스 추가
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        # 데이터 정리
        report_text = f"""
시스템 현황 ({datetime.now().strftime('%m월 %d일')})

HDD 사용량:
"""
        for hdd in hdd_results:
            report_text += f"  {hdd}\n"
        
        report_text += "\nCPU 사용량:\n"
        for cpu in cpu_results:
            report_text += f"  {cpu}\n"
        
        report_text += "\nMemory 사용량:\n"
        for memory in memory_results:
            report_text += f"  {memory}\n"
        
        report_text += "\n문서 수:\n"
        for docs in docs_results:
            report_text += f"  {docs}\n"
        
        text_frame.text = report_text
        
        # 폰트 설정
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = '맑은 고딕'
                run.font.size = Pt(12)
        
        # 파일 저장
        # 명령행 인수로 파일명이 전달되면 사용, 아니면 기본 파일명 사용
        if len(sys.argv) > 1:
            filename = sys.argv[1]
        else:
            current_date = datetime.now()
            year = current_date.year
            month = current_date.month
            filename = f"RnD_월간보고_{year}년_{month:02d}월_현황.pptx"
        
        output_path = os.path.join(script_dir, filename)
        prs.save(output_path)
        
        print(f"기본 PowerPoint 파일이 생성되었습니다: {output_path}")
        
    except Exception as e:
        print(f"기본 PowerPoint 파일 생성 중 오류: {e}")
        import traceback
        traceback.print_exc()

print("--------------------------------")
